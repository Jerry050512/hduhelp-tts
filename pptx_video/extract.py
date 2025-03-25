import os
from pptx import Presentation
import win32com.client

def extract_text_from_ppt(pptx_path, output_dir):
    # 创建输出目录（如果不存在）
    os.makedirs(output_dir, exist_ok=True)

    # 加载PPTX文件
    prs = Presentation(pptx_path)

    text_result = []

    # 遍历每张幻灯片
    for slide_idx, slide in enumerate(prs.slides, start=1):
        # 准备文本文件路径
        txt_path = os.path.join(output_dir, f"slide-{slide_idx:04d}.txt")

        # 提取文本
        text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
        text_content = "\n".join(text)

        text_result.append(f"[Slide-{slide_idx}]\n{text_content}")

        # 保存文本到文件
        with open(txt_path, "w", encoding="utf-8") as f:
            f.write(text_content)
        print(f"保存文本到: {txt_path}")
    
    return text_result

def export_slides_as_images(pptx_path, output_dir):
    # 创建输出目录（如果不存在）
    os.makedirs(output_dir, exist_ok=True)

    # 使用win32com.client启动PowerPoint
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    powerpoint.Visible = True  # 设置为True以便可以看到操作过程

    # 打开PPTX文件
    presentation = powerpoint.Presentations.Open(pptx_path)

    # 导出每张幻灯片为PNG图片
    for slide_idx in range(presentation.Slides.Count):
        # 准备图片文件路径
        png_path = os.path.join(output_dir, f"slide-{slide_idx+1:04d}.png")

        # 导出幻灯片
        presentation.Slides[slide_idx].Export(png_path, "PNG")
        print(f"导出图片到: {png_path}")

    # 关闭PowerPoint
    presentation.Close()
    powerpoint.Quit()

def main():
    # 设置输入PPTX文件路径和输出目录
    input_pptx = r"C:\Users\jerry\Documents\project\hduhelp-tts\assets\demo\TTS-Demo.pptx"  # 替换为你的PPTX文件路径
    output_directory = r"C:\Users\jerry\Documents\project\hduhelp-tts\assets\processed\ppt"  # 替换为你的输出目录

    # 清空输出目录
    if os.path.exists(output_directory):
        for file in os.listdir(output_directory):
            os.remove(os.path.join(output_directory, file))

    # 提取文本
    text_result = extract_text_from_ppt(input_pptx, output_directory)

    # 导出幻灯片为图片
    export_slides_as_images(input_pptx, output_directory)

    with open(os.path.join(output_directory, "text.txt"), "w", encoding="utf-8") as f:
        f.write("\n---\n".join(text_result))

    print("任务完成！")

if __name__ == "__main__":
    main()